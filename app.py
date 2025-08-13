import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime
import io
import tempfile

# --- Page config ---
st.set_page_config(page_title="📝 PPT Data Entry", layout="centered")

# --- Custom CSS for light blue background and black labels ---
st.markdown(
    """
    <style>
    .stApp {
        background-color: #d0e7ff;  /* Light blue background */
    }
    /* Make labels and expander titles black */
    label, .stExpander > div > div > div > span {
        color: black !important;
        font-weight: bold !important;
    }
    </style>
    """,
    unsafe_allow_html=True
)

st.markdown("<h1 style='text-align: center; color: darkblue;'>📊 PPT Data Entry Form</h1>", unsafe_allow_html=True)
st.write("---")

date_format = "%d-%m-%Y"

# --- FORM IN EXPANDERS ---
with st.form("data_form"):

    # --- Step 1: Basic Info ---
    with st.expander("1️⃣ Basic Information", expanded=True):
        st.markdown("**Plant Name**")
        plant_name = st.text_input("", placeholder="Enter Plant Name")
        st.markdown("**Equipment**")
        equipment_name = st.text_input("", placeholder="Enter Equipment Name")
        st.markdown("**Case Enabler**")
        case_enabler = st.text_input("", placeholder="Enter Case Enabler")
        st.markdown("**Downtime Hours**")
        downtime_hours = st.text_input("", placeholder="Enter downtime hours")
        st.markdown("**📷 Upload Equipment Image**")
        equipment_image = st.file_uploader("", type=['png','jpg','jpeg'])

    # --- Step 2: Observation & Recommendation ---
    with st.expander("2️⃣ Observation & Recommendation"):
        st.markdown("**Observation Date**")
        observation_date = st.date_input("")
        st.markdown("**Observation**")
        observation = st.text_input("", placeholder="Enter observation details")
        st.markdown("**Date of Recommendation**")
        date_recommendation = st.date_input("")
        st.markdown("**Recommendation**")
        recommendation = st.text_input("", placeholder="Enter recommendation")

    # --- Step 3: Corrective Actions ---
    with st.expander("3️⃣ Corrective Actions"):
        st.markdown("**Date of Corrective Action Taken**")
        date_corrective_action = st.date_input("")
        st.markdown("**Corrective Action Details**")
        corrective_action_details = st.text_input("")
        st.markdown("**Date of Closed Report**")
        date_closed_report = st.date_input("")
        st.markdown("**Closed Report Status**")
        closed_report_status = st.text_input("")

    # --- Step 4: Machine & Trends ---
    with st.expander("4️⃣ Machine Details & Trends"):
        st.markdown("**Machine Details**")
        machine_details = st.text_area("", placeholder="Enter machine details")
        st.markdown("**Trend for Image-1**")
        trend_for_image1 = st.text_input("", placeholder="Enter trend description")
        st.markdown("**📷 Upload Trend Image 1**")
        trend_image1 = st.file_uploader("", type=['png','jpg','jpeg'])
        st.markdown("**Trend for Image-2**")
        trend_for_image2 = st.text_input("", placeholder="Enter trend description")
        st.markdown("**📷 Upload Trend Image 2**")
        trend_image2 = st.file_uploader("", type=['png','jpg','jpeg'])

    submit = st.form_submit_button("✅ Generate PPT")

# --- ON SUBMIT ---
if submit:
    template_path = "template.pptx"
    prs = Presentation(template_path)
    
    replacements = {
        "{{Plant Name}}": plant_name,
        "{{Equipment}}": equipment_name,
        "{{Case Enabler}}": case_enabler,
        "{{Downtime Hours}}": downtime_hours,
        "{{Observation Date}}": observation_date.strftime(date_format),
        "{{Observation}}": observation,
        "{{Date of Recommendation}}": date_recommendation.strftime(date_format),
        "{{Recommendation}}": recommendation,
        "{{Date of Corrective action Taken}}": date_corrective_action.strftime(date_format),
        "{{Corrective Action Details}}": corrective_action_details,
        "{{Date of closed Report}}": date_closed_report.strftime(date_format),
        "{{Closed Report Status}}": closed_report_status,
        "{{Machine Details}}": machine_details,
        "{{Trend for Image-1}}": trend_for_image1,
        "{{Trend for Image-2}}": trend_for_image2
    }

    def replace_placeholder(para, replacements):
        full_text = "".join(run.text for run in para.runs)
        for key, value in replacements.items():
            full_text = full_text.replace(key, value)
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = full_text
        else:
            para.add_run().text = full_text

    # Replace text and insert images
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Equipment Image
                if "{{Equipment Image}}" in shape.text and equipment_image:
                    shape.text = ""
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(equipment_image.getvalue())
                        tmp_path = tmp.name
                    slide.shapes.add_picture(tmp_path, left, top, width, height)
                    os.unlink(tmp_path)
                # Trend Image 1
                elif "{{Trend Image 1}}" in shape.text and trend_image1:
                    shape.text = ""
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(trend_image1.getvalue())
                        tmp_path = tmp.name
                    slide.shapes.add_picture(tmp_path, left, top, width, height)
                    os.unlink(tmp_path)
                # Trend Image 2
                elif "{{Trend Image 2}}" in shape.text and trend_image2:
                    shape.text = ""
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        tmp.write(trend_image2.getvalue())
                        tmp_path = tmp.name
                    slide.shapes.add_picture(tmp_path, left, top, width, height)
                    os.unlink(tmp_path)
                else:
                    for para in shape.text_frame.paragraphs:
                        replace_placeholder(para, replacements)

    # Save PPT to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)

    st.success("🎉 PPT generated successfully!")
    st.download_button(
        "⬇️ Download PPT",
        ppt_bytes,
        file_name=f"{equipment_name}_{plant_name}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
